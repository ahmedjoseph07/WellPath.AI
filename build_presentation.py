"""Generate a detailed, slate-tinted light-theme thesis presentation.

Goals:
- Light theme with a darker flavour (slate-200 background, white cards,
  deep-teal accent) — same design sense as the original dark deck.
- Deeper algorithmic coverage: decision trees -> bagging vs boosting ->
  gradient boosting -> XGBoost innovations, plus per-operator GA detail
  and explicit MCM formulas.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Palette (darker-flavour light theme) ----------
BG        = RGBColor(0xE6, 0xEC, 0xF2)   # slate-200 tint
CARD      = RGBColor(0xFF, 0xFF, 0xFF)
SOFT      = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100 (zebra rows)
BORDER    = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
ACCENT    = RGBColor(0x0E, 0x74, 0x90)   # darker cyan
ACCENT_D  = RGBColor(0x0C, 0x4A, 0x6E)   # deeper for emphasis text
ACCENT_L  = RGBColor(0xEC, 0xFE, 0xFF)   # cyan-50 fill for callouts
ACCENT_LB = RGBColor(0x67, 0xE8, 0xF9)   # cyan-300 border for callouts
INK       = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
INK_2     = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
MUTED     = RGBColor(0x47, 0x55, 0x69)   # slate-600
FAINT     = RGBColor(0x64, 0x74, 0x8B)   # slate-500
OK        = RGBColor(0x15, 0x80, 0x3D)
WARN      = RGBColor(0xB4, 0x53, 0x09)
BAD       = RGBColor(0xB9, 0x1C, 0x1C)

FONT_SANS = "Inter"
FONT_MONO = "JetBrains Mono"

SW, SH = Emu(12191695), Emu(6858000)   # 16:9 widescreen

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
blank = prs.slide_layouts[6]


# ---------- primitives ----------
def rect(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def round_rect(slide, x, y, w, h, fill, line=None, line_w=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, body, *, size=12, bold=False, color=INK,
         font=FONT_SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         italic=False, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def chrome(slide, section_num=None, part_label=None):
    set_bg(slide, BG)
    rect(slide, 0, 0, SW, Emu(58000), ACCENT)
    if section_num is not None:
        round_rect(slide, Emu(480000), Emu(360000), Emu(520000), Emu(520000),
                   ACCENT, radius=0.5)
        text(slide, Emu(480000), Emu(360000), Emu(520000), Emu(520000),
             f"{section_num:02d}", size=18, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Emu(480000), Emu(6460000), SW - Emu(960000), Emu(280000),
         "WellPath.AI  ·  Joseph Ahmed (2007007)  ·  CUET  ·  Dec 2025",
         size=9, color=FAINT, align=PP_ALIGN.RIGHT)
    if part_label:
        text(slide, Emu(480000), Emu(6460000), SW - Emu(960000), Emu(280000),
             part_label, size=9, color=FAINT, align=PP_ALIGN.LEFT)


def title_block(slide, t, sub=None, x=Emu(1100000), y=Emu(380000),
                w=Emu(10500000)):
    text(slide, x, y, w, Emu(440000), t, size=26, bold=True, color=INK)
    if sub:
        text(slide, x, y + Emu(440000), w, Emu(280000),
             sub, size=12, color=MUTED)


def bullet(slide, x, y, w, t, *, size=12, color=INK, dot_color=ACCENT,
           dot_size=Emu(70000)):
    rect(slide, x, y + Emu(70000), dot_size, dot_size, dot_color)
    text(slide, x + Emu(180000), y, w - Emu(180000), Emu(260000),
         t, size=size, color=color)


def section_label(slide, x, y, w, t):
    text(slide, x, y, w, Emu(260000), t,
         size=10, bold=True, color=ACCENT_D)


def callout(slide, y, t, w=Emu(10400000), x=Emu(1100000), h=Emu(380000),
            size=11):
    round_rect(slide, x, y, w, h, ACCENT_L,
               line=ACCENT_LB, line_w=Emu(9525), radius=0.3)
    text(slide, x, y, w, h, t, size=size, bold=True, color=ACCENT_D,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG)
rect(s, 0, 0, Emu(120000), SH, ACCENT)
for i in range(6):
    rect(s, Emu(480000 + i*140000), Emu(540000),
         Emu(80000), Emu(80000),
         ACCENT if i % 2 == 0 else RGBColor(0x67, 0xE8, 0xF9))

text(s, Emu(480000), Emu(800000), Emu(10000000), Emu(380000),
     "THESIS DEFENSE  ·  BSc Petroleum & Mining Engineering",
     size=12, bold=True, color=ACCENT_D)

text(s, Emu(480000), Emu(1180000), Emu(10500000), Emu(900000),
     "WellPath.AI", size=72, bold=True, color=INK)

text(s, Emu(480000), Emu(2060000), Emu(10500000), Emu(560000),
     "AI-Assisted Well Path Optimization",
     size=26, color=INK_2)
text(s, Emu(480000), Emu(2400000), Emu(10500000), Emu(440000),
     "Integrating Geosteering Principles with Historical Well Log Data",
     size=18, color=MUTED, italic=True)

rect(s, Emu(480000), Emu(3200000), Emu(800000), Emu(40000), ACCENT)

text(s, Emu(480000), Emu(3320000), Emu(10500000), Emu(340000),
     "Joseph Ahmed   ·   Student ID 2007007",
     size=18, bold=True, color=INK)
text(s, Emu(480000), Emu(3640000), Emu(10500000), Emu(300000),
     "Supervisor: Assistant Professor Aqif Hosain Khan",
     size=14, color=INK_2)
text(s, Emu(480000), Emu(3920000), Emu(10500000), Emu(300000),
     "Department of Petroleum & Mining Engineering, CUET",
     size=14, color=INK_2)

text(s, Emu(480000), Emu(5900000), Emu(10500000), Emu(300000),
     "December 2025", size=12, color=MUTED)
text(s, Emu(480000), Emu(6200000), Emu(10500000), Emu(280000),
     "github.com/ahmedjoseph07/WellPath.AI",
     size=11, color=ACCENT_D, font=FONT_MONO)


# =====================================================================
# SLIDE 2 — Outline
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 1)
title_block(s, "Presentation Outline",
            "Detailed structure · ~35 min · 14 sections")

items = [
    ("Introduction & Problem",        "3"),
    ("Objectives & Scope",            "2"),
    ("Literature & Research Gap",     "3"),
    ("Petrophysics · Well Logs",      "3"),
    ("Decision Trees (foundation)",   "2"),
    ("Bagging vs Boosting",           "2"),
    ("XGBoost · Theory & Math",       "4"),
    ("XGBoost · Feature Importance",  "2"),
    ("Genetic Algorithm · Theory",    "3"),
    ("GA · Operators in Detail",      "4"),
    ("Minimum Curvature Method",      "2"),
    ("System & 3D Visualization",     "2"),
    ("Live Demo",                     "2"),
    ("Results · Limitations · Q&A",   "1"),
]
col_w = Emu(5200000); row_h = Emu(340000)
x0 = Emu(1100000); y0 = Emu(1300000)
for i, (name, t) in enumerate(items):
    col = i // 7; row = i % 7
    x = x0 + col * (col_w + Emu(400000))
    y = y0 + row * row_h
    text(s, x, y, Emu(480000), row_h, f"{i+1:02d}",
         size=15, bold=True, color=ACCENT_D, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Emu(540000), y, col_w - Emu(1200000), row_h,
         name, size=13, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    round_rect(s, x + col_w - Emu(620000), y + Emu(60000),
               Emu(600000), Emu(220000),
               ACCENT_L, line=ACCENT_LB, line_w=Emu(6350), radius=0.5)
    text(s, x + col_w - Emu(620000), y + Emu(60000),
         Emu(600000), Emu(220000),
         f"{t} min", size=9, bold=True, color=ACCENT_D,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 3 — Background & Motivation
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 2, "Part I · Introduction")
title_block(s, "Background & Motivation",
            "Why integrated ML-driven well planning is needed today")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "WHY WELL PATH PLANNING MATTERS")
why = [
    "Well placement is THE single biggest lever on productivity & recovery",
    "Traditional workflow: manual log picking → iterative trajectory design",
    "Requires tight loops between geology, petrophysics, drilling — days to weeks",
    "Commercial suites (Landmark WellPlan, Petrel) cost USD 50–200 k / year",
    "Regional gap: AI-assisted drilling in Bangladesh largely unexplored",
    "Opportunity: combine ML prediction + evolutionary search + 3D viz in one pipeline",
]
y = Emu(1580000)
for t in why:
    bullet(s, Emu(1100000), y, Emu(5000000), t, size=12)
    y += Emu(340000)

# right: stat cards
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "BY THE NUMBERS")
stats = [
    ("50 – 200 k",   "USD / year — commercial suite license"),
    ("3 – 7 days",   "Typical re-plan cycle for a directional well"),
    ("USD 6 – 15 M", "Average cost of a deviated exploration well"),
    ("±5 °",         "Inclination uncertainty between interpreters"),
    ("70 %+",        "Recovery variance traceable to trajectory design"),
]
y = Emu(1580000)
for v, d in stats:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(500000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    rect(s, Emu(6500000), y, Emu(80000), Emu(500000), ACCENT)
    text(s, Emu(6620000), y, Emu(2000000), Emu(500000),
         v, size=16, bold=True, color=INK, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8620000), y, Emu(2800000), Emu(500000),
         d, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(560000)


# =====================================================================
# SLIDE 4 — Problem Statement
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 3, "Part I · Introduction")
title_block(s, "Problem Statement",
            "Five concrete shortcomings of the current practice")

problems = [
    ("Subjectivity",
     "Well log interpretation varies between analysts",
     "→ Inconsistent reservoir picks and zone boundaries"),
    ("Disconnected tools",
     "Log analysis, trajectory design, visualization are separate manual steps",
     "→ Data re-entry, version drift, hand-offs across teams"),
    ("Time cost",
     "Multiple planning rounds between specialists — days to weeks per well",
     "→ Delays first oil, raises rig standby cost"),
    ("No mathematical optimum",
     "Trajectory design is driven by experience, not a fitness function",
     "→ Leaves productive zone contact on the table"),
    ("Cost barrier",
     "Commercial tools cost USD 50–200 k / year — out of reach for academia & small E&P",
     "→ Regional teams cannot iterate on modern methods"),
]
x0 = Emu(1100000); y = Emu(1300000)
w = Emu(10400000); h = Emu(920000); gap = Emu(60000)
for i, (n, d, impact) in enumerate(problems):
    round_rect(s, x0, y, w, h, CARD,
               line=BORDER, line_w=Emu(6350), radius=0.1)
    rect(s, x0, y, Emu(80000), h, ACCENT)
    round_rect(s, x0 + Emu(220000), y + Emu(240000),
               Emu(420000), Emu(420000), ACCENT, radius=0.5)
    text(s, x0 + Emu(220000), y + Emu(240000),
         Emu(420000), Emu(420000),
         f"{i+1}", size=16, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + Emu(740000), y + Emu(140000),
         Emu(2600000), Emu(280000), n,
         size=14, bold=True, color=INK)
    text(s, x0 + Emu(740000), y + Emu(440000),
         w - Emu(900000), Emu(260000), d,
         size=11, color=INK_2)
    text(s, x0 + Emu(740000), y + Emu(700000),
         w - Emu(900000), Emu(220000), impact,
         size=10, color=ACCENT_D, italic=True)
    y += h + gap


# =====================================================================
# SLIDE 5 — Research Objectives
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 4, "Part I · Introduction")
title_block(s, "Research Objectives",
            "Three deliverables that collectively form the WellPath.AI pipeline")

objs = [
    ("01", "Predict productive zones",
     "Train an XGBoost classifier on 5 historical well log curves "
     "(GR, Rt, RHOB, NPHI, DT) to assign each depth to one of three classes: "
     "productive, marginal, non-productive.",
     ["XGBClassifier (150 trees, max_depth=5)",
      "Heuristic pseudo-labels from petrophysical cutoffs",
      "Transductive: re-trained per well — no pretraining needed"]),
    ("02", "Optimize trajectory",
     "Use a Genetic Algorithm to evolve 3D well paths that maximize "
     "productive zone contact while respecting geological and drilling "
     "constraints (DLS limit, inclination bounds).",
     ["DEAP framework · population 50 · 100 generations",
      "Tournament selection + SBX crossover + polynomial mutation",
      "Fitness = zone exposure − 0.3 × max(DLS)"]),
    ("03", "Visualize & interact",
     "Deliver an interactive 3D web platform showing optimized well paths "
     "with formation layers, survey stations, DLS bands and rig geometry — "
     "decision-ready, in the browser.",
     ["React 18 + Vite + Three.js scene graph",
      "Formation boxes, KOP & TD markers, depth axis, compass",
      "Survey station coordinate table (exportable)"]),
]
x0 = Emu(1100000); y = Emu(1280000)
w = Emu(3400000); h = Emu(4700000); gap = Emu(240000)
for i, (n, t, body, feats) in enumerate(objs):
    x = x0 + i * (w + gap)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, y, Emu(80000), h, ACCENT)
    text(s, x + Emu(200000), y + Emu(200000),
         Emu(700000), Emu(340000),
         n, size=22, bold=True, color=ACCENT_D, font=FONT_MONO)
    text(s, x + Emu(200000), y + Emu(580000),
         w - Emu(300000), Emu(340000),
         t, size=15, bold=True, color=INK)
    text(s, x + Emu(200000), y + Emu(960000),
         w - Emu(300000), Emu(1800000),
         body, size=11, color=MUTED)
    fy = y + Emu(2850000)
    for ft in feats:
        bullet(s, x + Emu(200000), fy, w - Emu(300000), ft, size=10)
        fy += Emu(360000)


# =====================================================================
# SLIDE 6 — Scope (in / out)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 5, "Part I · Introduction")
title_block(s, "Scope of the Study",
            "What the thesis covers — and what is deliberately left for future work")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000), "IN SCOPE")
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000), "OUT OF SCOPE")

in_scope = [
    "Historical well log datasets (GR, Rt, RHOB, NPHI, DT)",
    "XGBoost classifier for reservoir zone productivity",
    "Genetic Algorithm trajectory optimization with DLS constraint",
    "Minimum Curvature Method for 3D coordinate conversion",
    "React + Three.js 3D visualization web platform",
    "Deterministic synthetic dataset (seed 2007007)",
    "Open-source release on GitHub — reproducible pipeline",
]
out_scope = [
    "Real-time MWD/LWD geosteering integration",
    "Multi-well correlation / geostatistical modelling",
    "Torque, drag, wellbore stability analysis",
    "Production simulation and economic evaluation",
    "Real Bangladesh field data (proof-of-concept only)",
    "Multi-lateral and fishbone trajectory designs",
    "Pore-pressure / fracture-gradient prediction",
]
y = Emu(1600000)
for a, b in zip(in_scope, out_scope):
    round_rect(s, Emu(1100000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, Emu(1100000), y, Emu(60000), Emu(360000), ACCENT)
    text(s, Emu(1220000), y, Emu(4800000), Emu(360000),
         a, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, Emu(6500000), y, Emu(60000), Emu(360000), BAD)
    text(s, Emu(6620000), y, Emu(4800000), Emu(360000),
         b, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(440000)


# =====================================================================
# SLIDE 7 — Literature Review
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 6, "Part I · Introduction")
title_block(s, "Literature Review",
            "Five representative prior works on ML and optimization for well planning")

papers = [
    ("Chen (2025)",
     "XGBoost for layer-wise production in commingled wells",
     "Showed XGBoost outperforms linear regression on multi-layer allocation. "
     "Motivates our choice of XGBoost for per-depth zone classification."),
    ("Timonov et al. · SPE-207364-MS",
     "ML-assisted geosteering with automated trajectory adjustment during drilling",
     "Closest operational analogue — but trajectory updates driven by real-time "
     "LWD, not historical logs. WellPath.AI is offline / pre-drill."),
    ("Mansouri et al. (2015)",
     "Multi-objective GA for 3D well path with DLS and target constraints",
     "Establishes GA as a viable trajectory optimizer. Our fitness function "
     "simplifies their Pareto setup into a single weighted scalar."),
    ("Zhou & Hu (2025)",
     "Multi-objective GA for clustered infill drilling",
     "Simultaneous multi-well optimization — extends Mansouri's framework. "
     "Out-of-scope here; single-well is our MVP."),
    ("AIP Conf. (2022)",
     "Review: Well placement optimization",
     "Confirms industry trend toward ML + optimization integration. "
     "Identifies the visualization-layer gap that our 3D viewer fills."),
]
x0 = Emu(1100000); y = Emu(1280000)
w = Emu(10400000); h = Emu(960000); gap = Emu(40000)
for au, t, why in papers:
    round_rect(s, x0, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.1)
    rect(s, x0, y, Emu(80000), h, ACCENT)
    text(s, x0 + Emu(220000), y + Emu(140000),
         Emu(4000000), Emu(240000),
         au, size=12, bold=True, color=ACCENT_D, font=FONT_MONO)
    text(s, x0 + Emu(220000), y + Emu(380000),
         w - Emu(400000), Emu(260000),
         t, size=12, bold=True, color=INK)
    text(s, x0 + Emu(220000), y + Emu(620000),
         w - Emu(400000), Emu(320000),
         why, size=10, color=MUTED)
    y += h + gap


# =====================================================================
# SLIDE 8 — Research Gaps
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 7, "Part I · Introduction")
title_block(s, "Research Gaps",
            "Four gaps in prior work — and what WellPath.AI contributes")

gaps = [
    ("No 3D visualization",
     "Most studies publish numerical outputs only",
     "We deliver an interactive Three.js scene with formations, trajectory, markers"),
    ("Regional gap",
     "Local Bangladesh datasets rarely used; AI-assisted drilling unexplored",
     "Thesis is seeded with regional context; architecture accepts real field data"),
    ("No ML + GA + 3D hybrid",
     "Literature lacks a single system spanning prediction → optimization → visualization",
     "WellPath.AI is exactly that end-to-end pipeline"),
    ("No geosteering in GA fitness",
     "Existing GA trajectory work does not use zone prediction to guide fitness",
     "Our fitness interpolates ML zone scores directly along the decoded path"),
]
x0 = Emu(1100000); y = Emu(1280000)
w = Emu(10400000); h = Emu(1080000); gap = Emu(40000)
for n, prev, ours in gaps:
    round_rect(s, x0, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.1)
    rect(s, x0, y, Emu(80000), h, ACCENT)
    text(s, x0 + Emu(220000), y + Emu(180000),
         w - Emu(400000), Emu(280000),
         n, size=14, bold=True, color=INK)
    text(s, x0 + Emu(220000), y + Emu(500000),
         Emu(4800000), Emu(520000),
         "PRIOR WORK", size=9, bold=True, color=MUTED)
    text(s, x0 + Emu(220000), y + Emu(720000),
         Emu(4800000), Emu(300000),
         prev, size=10, color=INK_2)
    text(s, x0 + Emu(5400000), y + Emu(500000),
         Emu(5000000), Emu(520000),
         "WELLPATH.AI", size=9, bold=True, color=ACCENT_D)
    text(s, x0 + Emu(5400000), y + Emu(720000),
         Emu(5000000), Emu(320000),
         ours, size=10, color=INK_2)
    y += h + gap


# =====================================================================
# SLIDE 9 — What are Well Logs
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 8, "Part II · Petrophysics")
title_block(s, "What are Well Logs?",
            "Continuous depth measurements of rock & fluid properties — raw ML input")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000), "DEFINITION")
for i, t in enumerate([
    "Tools lowered on wireline — record physical properties vs depth",
    "Primary input to reservoir characterization",
    "Convert raw signals → lithology, porosity, fluid type",
    "Sampling interval typically 0.15 – 0.5 m (we use 5 m bins)",
    "WellPath.AI uses 5 curves as ML features (listed →)",
]):
    bullet(s, Emu(1100000), Emu(1580000) + i * Emu(340000),
           Emu(5000000), t, size=12)

section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "FIVE FEATURES USED BY THE MODEL")
logs = [
    ("GR",   "Gamma Ray",        "Natural radioactivity · shale vs sand"),
    ("Rt",   "Resistivity",      "Only log that separates HC vs water"),
    ("RHOB", "Density",          "Compton scatter · bulk density"),
    ("NPHI", "Neutron porosity", "Hydrogen index · porosity"),
    ("DT",   "Sonic",            "P-wave Δt · rock competence"),
]
y = Emu(1580000)
for code, name, why in logs:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(380000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.2)
    round_rect(s, Emu(6600000), y + Emu(60000), Emu(640000), Emu(260000),
               ACCENT, radius=0.3)
    text(s, Emu(6600000), y + Emu(60000), Emu(640000), Emu(260000),
         code, size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(7300000), y, Emu(1800000), Emu(380000),
         name, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(9200000), y, Emu(2200000), Emu(380000),
         why, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(440000)

callout(s, Emu(5900000),
        "Together the 5 curves form a multi-dimensional signal that XGBoost can learn — more powerful than any single cutoff",
        size=11)


# =====================================================================
# SLIDE 10 — Log Curve Detail Table
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 9, "Part II · Petrophysics")
title_block(s, "Log Curves — Physics & Productive Cutoffs",
            "How each signal is measured · typical range · reading · cutoff used for labels")

cols = [("Curve", 1300000), ("Physical principle", 3200000),
        ("Range", 1650000), ("Reading", 1950000),
        ("Productive cutoff", 2300000)]
rows = [
    ("GR",   "K, Th, U radioactivity in shale",   "0 – 150 API",    "Low = sand · High = shale",      "< 60 API"),
    ("Rt",   "Electrical current resistance",     "0.2 – 2000 Ω·m", "High = HC · Low = water",         "> 15 Ω·m"),
    ("RHOB", "Compton γ-ray scattering",          "1.9 – 2.9 g/cc", "Low = porous · High = tight",    "< 2.35 g/cc"),
    ("NPHI", "Hydrogen index (thermal neutrons)", "0 – 0.45 frac",  "High = porous · Low = tight/gas", "> 0.12"),
    ("DT",   "P-wave transit time",               "40 – 140 μs/ft", "High = soft/porous · Low = tight","60 – 70 μs/ft"),
]
x0 = Emu(1100000); y = Emu(1300000); row_h = Emu(420000)
cx = x0
for name, w in cols:
    rect(s, cx, y, Emu(w), row_h, ACCENT)
    text(s, cx + Emu(140000), y, Emu(w - 240000), row_h,
         name, size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         anchor=MSO_ANCHOR.MIDDLE)
    cx += Emu(w)
y += row_h
for ri, row in enumerate(rows):
    cx = x0
    bgc = CARD if ri % 2 == 0 else SOFT
    for ci, ((_, w), val) in enumerate(zip(cols, row)):
        rect(s, cx, y, Emu(w), row_h, bgc, line=BORDER, line_w=Emu(3175))
        bold = ci in (0, 4)
        color = ACCENT_D if ci == 4 else (INK if ci == 0 else INK_2)
        text(s, cx + Emu(140000), y, Emu(w - 240000), row_h,
             val, size=10, bold=bold, color=color,
             font=FONT_MONO if ci == 0 else FONT_SANS,
             anchor=MSO_ANCHOR.MIDDLE)
        cx += Emu(w)
    y += row_h

section_label(s, Emu(1100000), y + Emu(240000), Emu(10400000),
              "HOW TO READ THIS TABLE")
text(s, Emu(1100000), y + Emu(560000), Emu(10400000), Emu(700000),
     "A depth interval is labelled PRODUCTIVE only when Rt, GR and NPHI all pass their cutoffs simultaneously.\n"
     "MARGINAL if Rt and GR pass a relaxed criterion but NPHI fails. Everything else is NON-PRODUCTIVE.",
     size=11, color=INK_2)


# =====================================================================
# SLIDE 11 — Zone Classification
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 10, "Part II · Petrophysics")
title_block(s, "Zone Classification — Pseudo-Labels",
            "No DST/MDT ground truth available · labels derived from petrophysical cutoffs")

zones = [
    (OK, "1", "Productive",
     "Rt > 15 Ω·m   AND   GR < 60 API   AND   NPHI > 0.12",
     "Clean sand bearing hydrocarbons with significant porosity. "
     "This is the class the GA fitness function rewards visiting."),
    (WARN, "2", "Marginal",
     "Rt > 8 Ω·m   AND   GR < 80 API   (partial criteria only)",
     "Silty sand or partially HC-bearing interval. Uncertain reservoir quality — "
     "contributes a reduced score to the trajectory."),
    (BAD, "0", "Non-productive",
     "Any interval failing the above — high GR, low Rt, or low NPHI",
     "Shale, tight rock, or water-saturated formation. "
     "Trajectory should avoid spending measured depth here."),
]
x0 = Emu(1100000); y = Emu(1300000)
w = Emu(10400000); h = Emu(1450000); gap = Emu(80000)
for c, cls, name, rule, desc in zones:
    round_rect(s, x0, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.1)
    rect(s, x0, y, Emu(120000), h, c)
    round_rect(s, x0 + Emu(240000), y + Emu(200000),
               Emu(620000), Emu(380000),
               c, radius=0.2)
    text(s, x0 + Emu(240000), y + Emu(200000),
         Emu(620000), Emu(380000),
         f"class {cls}", size=10, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + Emu(960000), y + Emu(180000),
         Emu(3000000), Emu(340000),
         name, size=18, bold=True, color=INK)
    round_rect(s, x0 + Emu(240000), y + Emu(680000),
               w - Emu(400000), Emu(320000),
               SOFT, line=BORDER, line_w=Emu(3175), radius=0.2)
    text(s, x0 + Emu(240000), y + Emu(680000),
         w - Emu(400000), Emu(320000),
         rule, size=11, bold=True, color=INK,
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + Emu(240000), y + Emu(1060000),
         w - Emu(400000), Emu(340000),
         desc, size=11, color=MUTED)
    y += h + gap


# =====================================================================
# SLIDE 12 — Decision Trees (foundation of XGBoost)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 11, "Part III · Machine Learning")
title_block(s, "Decision Trees — The Foundation",
            "Every XGBoost model is an ensemble of these — start here")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "HOW A DECISION TREE WORKS")
for i, t in enumerate([
    "Start with all training samples at the root node",
    "Pick the feature + threshold that best separates classes",
    "Split data into two child nodes — recurse",
    "Stop when: max_depth reached OR pure node OR min_samples",
    "Leaf node stores the predicted class / value",
    "Prediction: walk the tree top-down following splits",
]):
    round_rect(s, Emu(1100000), Emu(1580000) + i * Emu(300000),
               Emu(380000), Emu(260000), ACCENT, radius=0.3)
    text(s, Emu(1100000), Emu(1580000) + i * Emu(300000),
         Emu(380000), Emu(260000), f"{i+1}",
         size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(1560000), Emu(1580000) + i * Emu(300000),
         Emu(4500000), Emu(260000), t, size=11, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)

# splitting criteria card
section_label(s, Emu(1100000), Emu(3540000), Emu(5000000),
              "SPLIT CRITERION (Gini impurity)")
round_rect(s, Emu(1100000), Emu(3860000), Emu(5000000), Emu(400000),
           SOFT, line=BORDER, line_w=Emu(6350), radius=0.15)
text(s, Emu(1100000), Emu(3860000), Emu(5000000), Emu(400000),
     "G = 1 − Σ pᵢ²     (lower G = purer node)",
     size=12, bold=True, color=INK, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Emu(1100000), Emu(4340000), Emu(5000000), Emu(600000),
     "Best split = maximises (G_parent − weighted_avg(G_children)).\n"
     "XGBoost replaces Gini with a loss-function-aware gain.",
     size=10, color=MUTED)

# right: schematic
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "VISUAL SCHEMATIC")
# root
root_x = Emu(8600000); root_y = Emu(1600000); nw = Emu(1200000); nh = Emu(380000)
def node(x, y, t, color=ACCENT):
    round_rect(s, x, y, nw, nh, color, radius=0.3)
    text(s, x, y, nw, nh, t, size=10, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
node(root_x, root_y, "GR < 60?")
# level 2
node(Emu(7400000), Emu(2400000), "Rt > 15?")
node(Emu(9800000), Emu(2400000), "NON-PROD", color=BAD)
# level 3 left
node(Emu(6400000), Emu(3200000), "NPHI > .12?")
node(Emu(8400000), Emu(3200000), "MARGINAL", color=WARN)
# level 4 leftmost
node(Emu(5400000), Emu(4000000), "PRODUCTIVE", color=OK)
node(Emu(7400000), Emu(4000000), "MARGINAL", color=WARN)

# simple connecting lines
def line(x1, y1, x2, y2):
    ln = s.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = MUTED
    ln.line.width = Emu(9525)
line(root_x + nw//2, root_y + nh, Emu(7400000) + nw//2, Emu(2400000))
line(root_x + nw//2, root_y + nh, Emu(9800000) + nw//2, Emu(2400000))
line(Emu(7400000) + nw//2, Emu(2400000) + nh, Emu(6400000) + nw//2, Emu(3200000))
line(Emu(7400000) + nw//2, Emu(2400000) + nh, Emu(8400000) + nw//2, Emu(3200000))
line(Emu(6400000) + nw//2, Emu(3200000) + nh, Emu(5400000) + nw//2, Emu(4000000))
line(Emu(6400000) + nw//2, Emu(3200000) + nh, Emu(7400000) + nw//2, Emu(4000000))

text(s, Emu(6500000), Emu(4600000), Emu(5000000), Emu(600000),
     "A single tree is easy to interpret but fragile —\nsensitive to noise, weak on unseen data.",
     size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 13 — Bagging vs Boosting
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 12, "Part III · Machine Learning")
title_block(s, "Bagging vs Boosting — Ensemble Methods",
            "How we go from one weak tree to a strong ensemble")

methods = [
    ("Bagging  (Random Forest)",
     "PARALLEL",
     "Train many trees independently on random subsets of data. "
     "Final prediction = average vote.",
     ["Trees built in parallel",
      "Each on a bootstrap sample",
      "Reduces variance — less overfit",
      "Example: Random Forest (Breiman 2001)"],
     ACCENT),
    ("Boosting  (XGBoost, LightGBM)",
     "SEQUENTIAL",
     "Trees are built one at a time. Each new tree focuses on the "
     "residual errors of the previous ensemble.",
     ["Trees built sequentially",
      "Each corrects previous errors",
      "Reduces bias AND variance",
      "Examples: AdaBoost, GBDT, XGBoost"],
     ACCENT_D),
]
x0 = Emu(1100000); y = Emu(1280000)
w = Emu(5100000); h = Emu(4400000); gap = Emu(200000)
for i, (title, mode, body, feats, col) in enumerate(methods):
    x = x0 + i * (w + gap)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, y, w, Emu(100000), col)
    text(s, x + Emu(240000), y + Emu(240000),
         w - Emu(480000), Emu(320000),
         title, size=15, bold=True, color=INK)
    round_rect(s, x + Emu(240000), y + Emu(620000),
               Emu(1800000), Emu(260000),
               ACCENT_L, line=ACCENT_LB, line_w=Emu(6350), radius=0.5)
    text(s, x + Emu(240000), y + Emu(620000),
         Emu(1800000), Emu(260000),
         mode, size=9, bold=True, color=ACCENT_D,
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Emu(240000), y + Emu(1000000),
         w - Emu(480000), Emu(940000),
         body, size=11, color=MUTED)
    fy = y + Emu(2000000)
    for ft in feats:
        bullet(s, x + Emu(240000), fy, w - Emu(480000), ft, size=11)
        fy += Emu(380000)

callout(s, Emu(5840000),
        "WellPath.AI chooses BOOSTING — sequential error correction gives better accuracy on small tabular data (200 samples)",
        size=11)


# =====================================================================
# SLIDE 14 — Gradient Boosting Math
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 13, "Part III · Machine Learning")
title_block(s, "Gradient Boosting — Mathematical Formulation",
            "The algorithm XGBoost implements, expressed as an additive model")

# main equation
round_rect(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(600000),
           ACCENT_L, line=ACCENT_LB, line_w=Emu(9525), radius=0.15)
text(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(600000),
     "F_m(x)  =  F_{m−1}(x)  +  η · h_m(x)",
     size=22, bold=True, color=ACCENT_D, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# explanation
section_label(s, Emu(1100000), Emu(2080000), Emu(5000000),
              "TERM BY TERM")
defs = [
    ("F_m(x)",     "Ensemble prediction after m trees"),
    ("F_{m−1}(x)", "Previous ensemble (already built)"),
    ("η",          "Learning rate — shrinks each tree (we use 0.08)"),
    ("h_m(x)",     "New weak learner (tree) to be added"),
]
y = Emu(2400000)
for t, d in defs:
    round_rect(s, Emu(1100000), y, Emu(5000000), Emu(320000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    text(s, Emu(1220000), y, Emu(1600000), Emu(320000),
         t, size=12, bold=True, color=ACCENT_D, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(2900000), y, Emu(3100000), Emu(320000),
         d, size=10, color=INK_2, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(360000)

# right: fitting step
section_label(s, Emu(6500000), Emu(2080000), Emu(5000000),
              "HOW h_m(x) IS FIT")
for i, t in enumerate([
    "Compute residuals: r_i = −∂L/∂F(x_i)",
    "These are the negative gradients of the loss L",
    "Fit a new tree h_m to predict r_i",
    "Add tree to ensemble with learning rate η",
    "Repeat for m = 1, 2, …, M (we use M = 150)",
]):
    round_rect(s, Emu(6500000), Emu(2400000) + i * Emu(360000),
               Emu(360000), Emu(260000), ACCENT, radius=0.3)
    text(s, Emu(6500000), Emu(2400000) + i * Emu(360000),
         Emu(360000), Emu(260000), f"{i+1}",
         size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(6940000), Emu(2400000) + i * Emu(360000),
         Emu(4600000), Emu(260000), t,
         size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

callout(s, Emu(5020000),
        "Each tree learns from where the PREVIOUS ensemble was wrong — that is why boosting sequentially reduces bias",
        size=11)


# =====================================================================
# SLIDE 15 — XGBoost Innovations
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 14, "Part III · Machine Learning")
title_block(s, "XGBoost — What Makes It Special",
            "eXtreme Gradient Boosting (Chen & Guestrin, 2016) — four concrete improvements")

innovations = [
    ("2nd-order Taylor expansion",
     "Obj ≈ Σ [g_i · h_m(x_i)  +  ½ · H_i · h_m(x_i)²]",
     "Classic GBDT only uses gradient g. XGBoost also uses the Hessian H — "
     "the second derivative of the loss. Result: more precise split decisions."),
    ("Regularized objective",
     "Obj = Σ L(y_i, ŷ_i)  +  γ·T  +  ½λ · ‖w‖²",
     "Explicit penalty on tree complexity T and leaf weights w. "
     "γ discourages splitting; λ shrinks leaf weights. Controls overfit."),
    ("Histogram-based splits",
     "tree_method = 'hist'  ·  256 bins per feature",
     "Features are pre-binned into histograms. Candidate splits search over "
     "bin edges, not every value. Fast on laptops, memory-efficient."),
    ("Shrinkage & column sampling",
     "η = 0.08   ·   subsample = 0.8   ·   colsample = 0.8",
     "Each tree's contribution is shrunk by η. Every tree sees only 80 % of "
     "rows and features. Three independent regularizers on top of L1/L2."),
]
x0 = Emu(1100000); y = Emu(1280000)
w = Emu(5100000); h = Emu(2400000); gx = Emu(200000); gy = Emu(200000)
for i, (name, formula, body) in enumerate(innovations):
    col = i % 2; row = i // 2
    x = x0 + col * (w + gx)
    yy = y + row * (h + gy)
    round_rect(s, x, yy, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, yy, w, Emu(80000), ACCENT)
    text(s, x + Emu(240000), yy + Emu(200000),
         w - Emu(480000), Emu(320000),
         name, size=14, bold=True, color=INK)
    round_rect(s, x + Emu(240000), yy + Emu(600000),
               w - Emu(480000), Emu(420000),
               SOFT, line=BORDER, line_w=Emu(3175), radius=0.15)
    text(s, x + Emu(240000), yy + Emu(600000),
         w - Emu(480000), Emu(420000),
         formula, size=11, bold=True, color=INK,
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Emu(240000), yy + Emu(1100000),
         w - Emu(480000), Emu(1200000),
         body, size=11, color=MUTED)


# =====================================================================
# SLIDE 16 — XGBoost Configuration & Why Over Alternatives
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 15, "Part III · Machine Learning")
title_block(s, "Model Configuration & Why XGBoost",
            "Hyperparameters we use · and why alternatives were rejected")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "MODEL HYPERPARAMETERS")
params = [
    ("n_estimators",      "150",  "number of trees"),
    ("max_depth",         "5",    "tree complexity"),
    ("learning_rate",     "0.08", "shrinkage per tree"),
    ("subsample",         "0.8",  "row sampling"),
    ("colsample_bytree",  "0.8",  "feature sampling"),
    ("gamma",             "0.1",  "min split gain"),
    ("reg_alpha  (L1)",   "0.1",  "Lasso penalty"),
    ("reg_lambda (L2)",   "1.0",  "Ridge penalty"),
    ("objective",         "multi:softprob", "3-class probabilistic"),
]
y = Emu(1580000)
for name, val, desc in params:
    round_rect(s, Emu(1100000), y, Emu(5000000), Emu(320000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    text(s, Emu(1220000), y, Emu(2100000), Emu(320000),
         name, size=10, color=INK, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(3320000), y, Emu(1400000), Emu(320000),
         val, size=11, bold=True, color=ACCENT_D, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(4720000), y, Emu(1300000), Emu(320000),
         desc, size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(360000)

section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "WHY NOT THE ALTERNATIVES?")
alts = [
    ("Logistic Regression", "Cannot capture non-linear log interactions",     BAD),
    ("Random Forest",       "No sequential correction → more trees needed",  BAD),
    ("SVM",                 "Kernel choice brittle, slower, no native FI",   BAD),
    ("Neural Network",      "Needs 1000s of samples; black-box features",   BAD),
    ("LightGBM",            "Similar perf but less mature ecosystem",        WARN),
    ("XGBoost",             "Fast · regularized · interpretable · proven",   OK),
]
y = Emu(1580000)
for name, reason, mark in alts:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    rect(s, Emu(6500000), y, Emu(60000), Emu(360000), mark)
    text(s, Emu(6620000), y, Emu(1800000), Emu(360000),
         name, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8420000), y, Emu(3000000), Emu(360000),
         reason, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(400000)

text(s, Emu(6500000), Emu(4180000), Emu(5000000), Emu(400000),
     "Ref: Shwartz-Ziv & Armon (2022, NeurIPS) — XGBoost beats neural nets on structured tabular data.",
     size=9, color=MUTED, italic=True)


# =====================================================================
# SLIDE 17 — Feature Importance
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 16, "Part III · Machine Learning")
title_block(s, "Feature Importance",
            "Which logs drive the classifier — and whether the ordering is physically correct")

feats = [
    ("Resistivity (Rt)",    35, "Strongest HC indicator — only log that separates HC vs water"),
    ("Gamma Ray (GR)",      28, "Lithology discriminator — shale vs reservoir rock"),
    ("Neutron Porosity",    18, "Confirms porosity — essential for productive classification"),
    ("Density (RHOB)",      12, "Complements neutron in a crossplot — refines lithology"),
    ("Sonic (DT)",           7, "Weakest separator in clastic settings — confirmatory role"),
]
x_label = Emu(1100000); x_bar = Emu(3500000); bar_max = Emu(6800000)
y = Emu(1400000); row_h = Emu(560000)
for name, pct, why in feats:
    text(s, x_label, y, Emu(2300000), Emu(260000),
         name, size=13, bold=True, color=INK)
    text(s, x_label, y + Emu(260000), Emu(2300000), Emu(240000),
         why, size=10, color=MUTED)
    round_rect(s, x_bar, y + Emu(120000), bar_max, Emu(300000),
               SOFT, line=BORDER, line_w=Emu(3175), radius=0.5)
    w = int(bar_max * pct / 40)
    round_rect(s, x_bar, y + Emu(120000), w, Emu(300000),
               ACCENT, radius=0.5)
    text(s, x_bar + w + Emu(100000), y + Emu(120000),
         Emu(800000), Emu(300000),
         f"{pct}%", size=12, bold=True, color=ACCENT_D,
         font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

callout(s, Emu(5900000),
        "Rt > GR > NPHI > RHOB > DT is the accepted ordering for clastic reservoirs — the model learned the right signal",
        size=11)


# =====================================================================
# SLIDE 18 — Training Pipeline
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 17, "Part III · Machine Learning")
title_block(s, "Training Pipeline",
            "Transductive — model is calibrated per-well, no pretraining required")

steps = [
    ("1", "Input",     "Log dict:\n{depths,\nGR, Rt, RHOB,\nNPHI, DT}"),
    ("2", "Normalize", "StandardScaler\non 5 features\n(μ=0, σ=1)"),
    ("3", "Label",     "Heuristic rules\n→ 3 classes\n0 · 1 · 2"),
    ("4", "Train",     "XGBClassifier.fit()\n150 trees\nmlogloss"),
    ("5", "Predict",   "predict_proba()\n→ P(productive)\n+ zone_label"),
]
x0 = Emu(1100000); y = Emu(1600000)
box_w = Emu(1900000); box_h = Emu(1900000); gap = Emu(200000)
cur = x0
for i, (n, name, body) in enumerate(steps):
    round_rect(s, cur, y, box_w, box_h, CARD,
               line=BORDER, line_w=Emu(6350), radius=0.15)
    round_rect(s, cur + Emu(100000), y + Emu(100000),
               Emu(480000), Emu(360000), ACCENT, radius=0.3)
    text(s, cur + Emu(100000), y + Emu(100000),
         Emu(480000), Emu(360000),
         n, size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cur + Emu(100000), y + Emu(560000),
         box_w - Emu(200000), Emu(300000),
         name, size=13, bold=True, color=INK)
    text(s, cur + Emu(100000), y + Emu(880000),
         box_w - Emu(200000), Emu(900000),
         body, size=10, color=MUTED, font=FONT_MONO)
    cur += box_w
    if i < len(steps) - 1:
        text(s, cur, y + (box_h // 2) - Emu(150000),
             gap, Emu(300000),
             "→", size=22, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cur += gap

callout(s, Emu(4200000),
        "Transductive = no held-out test set. The entire well is labelled, the model memorises the signal, prediction is then stable per-well.",
        size=11, w=Emu(10400000))


# =====================================================================
# SLIDE 19 — GA Theory
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 18, "Part IV · Genetic Algorithm")
title_block(s, "Genetic Algorithm — Evolutionary Theory",
            "Darwinian natural selection, expressed as a search algorithm")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "CORE IDEA")
for i, t in enumerate([
    "Maintain a population of candidate solutions",
    "Score each with a fitness function",
    "Best individuals reproduce — worst die off",
    "Offspring inherit traits + random variation",
    "Repeat for many generations → converge to optimum",
    "Gradient-free — fitness need not be differentiable",
]):
    bullet(s, Emu(1100000), Emu(1580000) + i * Emu(300000),
           Emu(5000000), t, size=11)

# right: biology ↔ algorithm
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "BIOLOGY  ↔  ALGORITHM")
mapping = [
    ("Chromosome",  "14-gene vector [Inc, Az] × 7 stations"),
    ("Gene",        "One inclination or azimuth angle"),
    ("Population",  "50 candidate trajectories"),
    ("Generation",  "One select–cross–mutate iteration"),
    ("Fitness",     "Zone exposure − 0.3 × max(DLS)"),
    ("Selection",   "Tournament of size 3"),
    ("Crossover",   "SBX  (η = 20, p = 0.7)"),
    ("Mutation",    "Polynomial bounded  (η = 20, p = 0.2/gene)"),
]
y = Emu(1580000)
for a, b in mapping:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(320000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    text(s, Emu(6620000), y, Emu(1700000), Emu(320000),
         a, size=11, bold=True, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8320000), y, Emu(3100000), Emu(320000),
         b, size=10, color=ACCENT_D, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(360000)


# =====================================================================
# SLIDE 20 — Chromosome Encoding & Decoding
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 19, "Part IV · Genetic Algorithm")
title_block(s, "Chromosome Encoding & Decoding",
            "How a well trajectory becomes a GA individual — and how it turns back into 3D")

# chromosome banner
round_rect(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(560000),
           ACCENT_L, line=ACCENT_LB, line_w=Emu(9525), radius=0.15)
text(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(560000),
     "Individual  =  [ Inc₁, Az₁, Inc₂, Az₂, Inc₃, Az₃, …, Inc₇, Az₇ ]   →   14 real-valued genes",
     size=15, bold=True, color=ACCENT_D, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

section_label(s, Emu(1100000), Emu(2060000), Emu(5000000), "ENCODING")
for i, t in enumerate([
    "8 survey stations along measured depth",
    "Station 0 (surface) is fixed: Inc = 0°, Az = 0°",
    "Remaining 7 stations → 2 genes each = 14 genes",
    "Inclination bounded: 0° – 85°",
    "Azimuth bounded: 0° – 360°",
    "Real-valued (float) encoding — not binary strings",
]):
    bullet(s, Emu(1100000), Emu(2400000) + i * Emu(320000),
           Emu(5000000), t, size=11)

section_label(s, Emu(6500000), Emu(2060000), Emu(5000000), "DECODING STEPS")
for i, t in enumerate([
    "Split chromosome into 7 (Inc, Az) pairs",
    "Space 8 survey stations evenly along MD",
    "Apply Minimum Curvature Method",
    "(MD, Inc, Az) → (N, E, TVD) per station",
    "Interpolate productivity score at each TVD",
    "Compute fitness from exposure − DLS penalty",
]):
    round_rect(s, Emu(6500000), Emu(2400000) + i * Emu(320000),
               Emu(360000), Emu(260000), ACCENT, radius=0.3)
    text(s, Emu(6500000), Emu(2400000) + i * Emu(320000),
         Emu(360000), Emu(260000), f"{i+1}",
         size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(6940000), Emu(2400000) + i * Emu(320000),
         Emu(4500000), Emu(260000), t,
         size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 21 — Fitness Function
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 20, "Part IV · Genetic Algorithm")
title_block(s, "Fitness Function",
            "What the GA is maximizing — productive contact traded against drilling risk")

round_rect(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(620000),
           ACCENT_L, line=ACCENT_LB, line_w=Emu(9525), radius=0.15)
text(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(620000),
     "F  =  P_exposure   −   w_DLS  ×  DLS_penalty",
     size=22, bold=True, color=ACCENT_D, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

terms = [
    ("P_exposure",  "Zone exposure",
     "Average productivity score along the decoded path.",
     "= (1/N) · Σ score(TVD_j)",
     "Range 0.0 (all non-prod) → 1.0 (all productive)"),
    ("DLS_penalty", "Drilling-safety penalty",
     "Maximum dogleg severity across all station pairs.",
     "DLS = (α / ΔMD) × 30   [°/30 m]",
     "Curbs sharp bends that risk pipe and casing"),
    ("w_DLS = 0.3", "Trade-off weight",
     "User-configurable scalar multiplier.",
     "Higher → smoother, less reservoir contact",
     "Lower → aggressive, more reservoir contact"),
]
x0 = Emu(1100000); y = Emu(2140000)
w = Emu(3400000); h = Emu(3400000); gap = Emu(240000)
for i, (t1, name, t2, t3, t4) in enumerate(terms):
    x = x0 + i * (w + gap)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, y, w, Emu(100000), ACCENT)
    text(s, x + Emu(200000), y + Emu(220000), w - Emu(400000),
         Emu(360000), t1, size=15, bold=True, color=ACCENT_D,
         font=FONT_MONO)
    text(s, x + Emu(200000), y + Emu(620000), w - Emu(400000),
         Emu(320000), name, size=13, bold=True, color=INK)
    text(s, x + Emu(200000), y + Emu(980000), w - Emu(400000),
         Emu(600000), t2, size=11, color=INK_2)
    round_rect(s, x + Emu(200000), y + Emu(1720000),
               w - Emu(400000), Emu(360000),
               SOFT, line=BORDER, line_w=Emu(3175), radius=0.15)
    text(s, x + Emu(200000), y + Emu(1720000),
         w - Emu(400000), Emu(360000),
         t3, size=10, bold=True, color=INK, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Emu(200000), y + Emu(2140000), w - Emu(400000),
         Emu(1000000), t4, size=10, color=MUTED)


# =====================================================================
# SLIDE 22 — Tournament Selection
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 21, "Part IV · GA Operators")
title_block(s, "Tournament Selection",
            "How parents are chosen — balance between pressure and diversity")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "ALGORITHM (size = 3)")
for i, t in enumerate([
    "Randomly pick 3 individuals from the population",
    "Compare their fitness values",
    "The fittest of the three wins → enters the mating pool",
    "Repeat until the mating pool is full (size = 50)",
    "Two winners form one breeding pair for crossover",
]):
    round_rect(s, Emu(1100000), Emu(1580000) + i * Emu(380000),
               Emu(380000), Emu(280000), ACCENT, radius=0.3)
    text(s, Emu(1100000), Emu(1580000) + i * Emu(380000),
         Emu(380000), Emu(280000), f"{i+1}",
         size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(1560000), Emu(1580000) + i * Emu(380000),
         Emu(4500000), Emu(280000), t,
         size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# right: why tournament
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "WHY NOT ROULETTE / RANK?")
cards = [
    ("Size = 1",  "No selection pressure — random walk"),
    ("Size = 2",  "Mild pressure — preserves diversity (we sometimes use this)"),
    ("Size = 3",  "Balanced — our chosen setting"),
    ("Size = 5+", "High pressure — fast convergence but risks premature stagnation"),
    ("Roulette",  "Fitness-proportional — breaks when fitness can be negative"),
    ("Rank",      "Robust but slow — requires sorting every generation"),
]
y = Emu(1580000)
for name, desc in cards:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    text(s, Emu(6620000), y, Emu(1400000), Emu(360000),
         name, size=11, bold=True, color=ACCENT_D, font=FONT_MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8000000), y, Emu(3400000), Emu(360000),
         desc, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(400000)


# =====================================================================
# SLIDE 23 — SBX Crossover
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 22, "Part IV · GA Operators")
title_block(s, "Simulated Binary Crossover (SBX)",
            "Combine two parent trajectories into offspring · real-valued, bounded")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "FORMULATION")
round_rect(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(700000),
           SOFT, line=BORDER, line_w=Emu(6350), radius=0.15)
text(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(700000),
     "c₁ = 0.5 [ (p₁+p₂) − β·(p₂−p₁) ]\n"
     "c₂ = 0.5 [ (p₁+p₂) + β·(p₂−p₁) ]",
     size=12, bold=True, color=INK, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for i, t in enumerate([
    "p₁, p₂ : parent gene values",
    "c₁, c₂ : offspring gene values",
    "β drawn from a polynomial distribution",
    "Controlled by η (eta) — distribution index",
    "High η → offspring close to parents",
    "Low η  → offspring spread far from parents",
]):
    bullet(s, Emu(1100000), Emu(2420000) + i * Emu(300000),
           Emu(5000000), t, size=11)

# config card
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "OUR CONFIGURATION")
cfg = [
    ("η (eta)",            "20",   "Tight distribution — children near parents"),
    ("Crossover prob.",    "0.7",  "70% of pairs crossover, 30% copy as-is"),
    ("Per-gene prob.",     "0.5",  "Each gene has 50% chance to be swapped"),
    ("Bounds",             "honored", "Offspring clipped to [low, high] per gene"),
    ("Framework",          "DEAP", "tools.cxSimulatedBinaryBounded"),
]
y = Emu(1580000)
for k, v, d in cfg:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    text(s, Emu(6620000), y, Emu(1500000), Emu(360000),
         k, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8120000), y, Emu(900000), Emu(360000),
         v, size=11, bold=True, color=ACCENT_D, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(9020000), y, Emu(2400000), Emu(360000),
         d, size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(400000)

callout(s, Emu(5140000),
        "SBX was designed specifically to mimic binary crossover's behaviour for real-valued chromosomes — ideal for [Inc, Az] genes",
        size=10)


# =====================================================================
# SLIDE 24 — Polynomial Mutation
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 23, "Part IV · GA Operators")
title_block(s, "Polynomial Mutation",
            "Random perturbation per gene · maintains late-stage exploration")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "FORMULATION")
round_rect(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(440000),
           SOFT, line=BORDER, line_w=Emu(6350), radius=0.15)
text(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(440000),
     "c = p + (high − low) · δ",
     size=13, bold=True, color=INK, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for i, t in enumerate([
    "p : original gene value",
    "c : mutated gene value",
    "[low, high] : bounds for this gene",
    "δ : polynomial-distributed perturbation",
    "Controlled by η — distribution index",
    "Low η → large mutations possible",
    "High η → small mutations preferred",
]):
    bullet(s, Emu(1100000), Emu(2180000) + i * Emu(300000),
           Emu(5000000), t, size=11)

# right: config + role
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "OUR CONFIGURATION & ROLE")
cfg2 = [
    ("η (eta)",         "20",   "Tight — small perturbations dominate"),
    ("Per-individual",  "1.0",  "Every individual is visited"),
    ("Per-gene prob.",  "0.2",  "≈ 3 of 14 genes mutate per chromosome"),
    ("Bounds",          "honored", "Clipped to [0–85°] and [0–360°]"),
    ("Framework",       "DEAP", "tools.mutPolynomialBounded"),
]
y = Emu(1580000)
for k, v, d in cfg2:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.12)
    text(s, Emu(6620000), y, Emu(1500000), Emu(360000),
         k, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8120000), y, Emu(900000), Emu(360000),
         v, size=11, bold=True, color=ACCENT_D, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(9020000), y, Emu(2400000), Emu(360000),
         d, size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(400000)

callout(s, Emu(5140000),
        "Without mutation the population collapses to duplicates of the best parent — mutation is what keeps the search alive",
        size=10)


# =====================================================================
# SLIDE 25 — GA Main Loop & Convergence
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 24, "Part IV · GA Operators")
title_block(s, "GA Main Loop & Convergence",
            "Putting the operators together — and how fitness evolves over 100 generations")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "MAIN LOOP (pseudocode)")
pseudo = [
    "initialise P₀ of size 50 (random within bounds)",
    "evaluate fitness of every individual",
    "for g in 1..100:",
    "    parents  = tournament(P, size=3)",
    "    children = SBX(parents, η=20, p=0.7)",
    "    children = polyMutate(children, η=20, p=0.2)",
    "    evaluate(children)",
    "    P = elitism(P ∪ children, keep=50)",
    "    record best fitness of P",
    "return argmax(P)",
]
round_rect(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(3500000),
           SOFT, line=BORDER, line_w=Emu(6350), radius=0.1)
ty = Emu(1680000)
for line in pseudo:
    text(s, Emu(1220000), ty, Emu(4800000), Emu(300000),
         line, size=11, color=INK, font=FONT_MONO)
    ty += Emu(320000)

# right: convergence phases
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "CONVERGENCE PHASES")
phases = [
    ("gen  1 – 20",  "Exploration",  "High diversity · random trajectories probed",  ACCENT),
    ("gen 20 – 70",  "Exploitation", "SBX combines good solutions · fast gains",      ACCENT),
    ("gen 70 – 100", "Convergence",  "Population homogeneous · marginal gains",       ACCENT_D),
]
y = Emu(1580000)
for rng, name, body, col in phases:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(820000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, Emu(6500000), y, Emu(80000), Emu(820000), col)
    text(s, Emu(6620000), y + Emu(100000),
         Emu(4800000), Emu(220000),
         rng, size=10, bold=True, color=ACCENT_D, font=FONT_MONO)
    text(s, Emu(6620000), y + Emu(320000),
         Emu(4800000), Emu(260000),
         name, size=13, bold=True, color=INK)
    text(s, Emu(6620000), y + Emu(560000),
         Emu(4800000), Emu(240000),
         body, size=10, color=MUTED)
    y += Emu(900000)

callout(s, Emu(5300000),
        "Typical gain 15 – 20 % from gen 1 → 100. Stochasticity is a feature — each run explores differently.",
        size=11)


# =====================================================================
# SLIDE 26 — Minimum Curvature Method (theory)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 25, "Part V · Geometry")
title_block(s, "Minimum Curvature Method (MCM)",
            "Industry-standard survey → 3D coordinate conversion (API / ISO)")

section_label(s, Emu(1100000), Emu(1250000), Emu(5000000),
              "ALGORITHM (between station i and i+1)")
round_rect(s, Emu(1100000), Emu(1580000), Emu(5000000), Emu(3700000),
           CARD, line=BORDER, line_w=Emu(6350), radius=0.1)
eqs = [
    "α  =  arccos[ cos(I₂−I₁) − sinI₁·sinI₂·(1−cos(A₂−A₁)) ]",
    "",
    "RF =  (2 / α) · tan(α/2)",
    "",
    "ΔN   = (ΔMD/2)(sinI₁·cosA₁ + sinI₂·cosA₂) · RF",
    "ΔE   = (ΔMD/2)(sinI₁·sinA₁ + sinI₂·sinA₂) · RF",
    "ΔTVD = (ΔMD/2)(cosI₁ + cosI₂) · RF",
    "",
    "When α → 0:  RF → 1   (L'Hôpital's rule)",
    "guarded numerically at α < 0.001 rad",
]
ty = Emu(1720000)
for line in eqs:
    text(s, Emu(1220000), ty, Emu(4800000), Emu(260000),
         line, size=11, color=INK, font=FONT_MONO)
    ty += Emu(300000)

# right: notes
section_label(s, Emu(6500000), Emu(1250000), Emu(5000000),
              "IMPLEMENTATION NOTES")
notes = [
    ("Assumption", "Wellbore follows a circular arc between stations"),
    ("Used by",    "WellPlan · COMPASS · WellPath.AI"),
    ("Frame",      "NED — North, East, Down"),
    ("Three.js",   "x = East  ·  y = −TVD  ·  z = North"),
    ("Accuracy",   "±1–2 m per 30 m survey interval"),
    ("Output",     "Array of {N, E, TVD} → tube geometry"),
    ("Why MCM",    "More accurate than tangential or balanced-tangential"),
]
y = Emu(1580000)
for k, v in notes:
    round_rect(s, Emu(6500000), y, Emu(5000000), Emu(360000),
               CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    text(s, Emu(6620000), y, Emu(1400000), Emu(360000),
         k, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(8020000), y, Emu(3400000), Emu(360000),
         v, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    y += Emu(400000)


# =====================================================================
# SLIDE 27 — Dogleg Severity
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 26, "Part V · Geometry")
title_block(s, "Dogleg Severity (DLS)",
            "Rate of wellbore curvature — why the fitness function penalizes it")

round_rect(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(540000),
           ACCENT_L, line=ACCENT_LB, line_w=Emu(9525), radius=0.15)
text(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(540000),
     "DLS  =  (α / ΔMD) × 30        [°/30 m]",
     size=22, bold=True, color=ACCENT_D, font=FONT_MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

probs = [
    ("Drill-pipe fatigue",
     "Repeated bending stress as pipe rotates through curved sections",
     "→ Fatigue cracking, twist-off, rig NPT"),
    ("Casing wear",
     "Drillstring contacts formation wall at high-curvature points",
     "→ Wall erosion, casing failure, costly workovers"),
    ("Torque & drag",
     "Increased friction along the string in curved sections",
     "→ Harder to reach TD, stuck-pipe risk"),
]
x0 = Emu(1100000); y = Emu(2040000)
w = Emu(3400000); h = Emu(1400000); gap = Emu(240000)
for i, (n, d, impact) in enumerate(probs):
    x = x0 + i * (w + gap)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, y, w, Emu(80000), ACCENT)
    text(s, x + Emu(200000), y + Emu(220000), w - Emu(400000),
         Emu(280000), n, size=13, bold=True, color=INK)
    text(s, x + Emu(200000), y + Emu(540000), w - Emu(400000),
         Emu(520000), d, size=10, color=MUTED)
    text(s, x + Emu(200000), y + Emu(1060000), w - Emu(400000),
         Emu(300000), impact, size=10, color=ACCENT_D, italic=True)

section_label(s, Emu(1100000), Emu(3700000), Emu(10400000),
              "SEVERITY BUCKETS (industry rule-of-thumb)")
buckets = [
    ("0 – 3°/30m", "Acceptable",    OK),
    ("3 – 5°/30m", "Moderate",      WARN),
    ("5 – 8°/30m", "Concerning",    RGBColor(0xEA, 0x58, 0x0C)),
    (">  8°/30m",  "Severe",        BAD),
]
bw = Emu(2520000); bx = Emu(1100000); by = Emu(4040000)
for rng, name, c in buckets:
    round_rect(s, bx, by, bw, Emu(560000), CARD,
               line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, bx, by, Emu(80000), Emu(560000), c)
    text(s, bx + Emu(180000), by + Emu(80000), bw - Emu(260000),
         Emu(240000), rng, size=11, bold=True, color=INK,
         font=FONT_MONO)
    text(s, bx + Emu(180000), by + Emu(300000), bw - Emu(260000),
         Emu(220000), name, size=10, color=MUTED)
    bx += bw + Emu(120000)

callout(s, Emu(5900000),
        "GA fitness subtracts w_DLS × max(DLS) — keeps the trajectory drillable, not just productive",
        size=11)


# =====================================================================
# SLIDE 28 — System Architecture
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 27, "Part VI · System")
title_block(s, "System Architecture",
            "Client-server via REST API · each layer independently testable")

layers = [
    ("Frontend",   "React 18 · Vite",
     ["Step-by-step workflow UI",
      "Three.js 3D scene graph",
      "Recharts 2D log plots",
      "Zustand state · Tailwind UI"]),
    ("REST API",   "FastAPI · Python",
     ["POST /api/upload",
      "POST /api/predict",
      "POST /api/optimize",
      "GET  /api/synthetic"]),
    ("ML Engine",  "XGBoost",
     ["StandardScaler normalize",
      "Heuristic pseudo-labels",
      "XGBClassifier · 150 trees",
      "Feature importance export"]),
    ("Optimizer",  "DEAP GA",
     ["Population init",
      "Tournament select",
      "SBX + polynomial mutate",
      "MCM decode · fitness eval"]),
]
x0 = Emu(1100000); y = Emu(1400000)
w = Emu(2500000); h = Emu(3300000); gap = Emu(120000)
for i, (name, tech, items) in enumerate(layers):
    x = x0 + i * (w + gap)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350))
    rect(s, x, y, w, Emu(80000), ACCENT)
    text(s, x + Emu(200000), y + Emu(200000), w - Emu(400000),
         Emu(320000), name, size=14, bold=True, color=INK)
    text(s, x + Emu(200000), y + Emu(540000), w - Emu(400000),
         Emu(280000), tech, size=10, color=ACCENT_D, font=FONT_MONO)
    ny = y + Emu(940000)
    for t in items:
        bullet(s, x + Emu(200000), ny, w - Emu(400000), t, size=10)
        ny += Emu(340000)
    if i < len(layers) - 1:
        text(s, x + w, y + (h // 2) - Emu(120000),
             gap, Emu(240000), "→", size=18, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

callout(s, Emu(5900000),
        "All four layers are open-source · stack runs on a laptop", size=11)


# =====================================================================
# SLIDE 29 — 3D Visualization Features
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 28, "Part VI · System")
title_block(s, "3D Visualization Features",
            "Interactive Three.js scene — decision-ready output in any browser")

viz = [
    ("Ground plane",      "Slab at TVD=0, 200 m scale bar"),
    ("Surface rig",       "Platform, derrick, BOP stack, wellhead"),
    ("Compass",           "N/S/E/W arrows with cone tips"),
    ("Depth axis",        "Tick labels — 'Ground Level 0m'"),
    ("Formation layers",  "Translucent — green / amber / red zones"),
    ("Well trajectory",   "Tube geometry following MCM survey"),
    ("Waypoint spheres",  "Colored by inclination"),
    ("KOP marker",        "Orange sphere + label (Inc > 3°)"),
    ("TD marker",         "Glowing cyan sphere at total depth"),
    ("Convergence chart", "GA fitness sparkline + % improvement"),
    ("Control modes",     "Orbit · Grab · Zoom (cursor feedback)"),
    ("Camera presets",    "Perspective · Front · Side · Top-down"),
]
x0 = Emu(1100000); y0 = Emu(1300000)
w = Emu(5100000); h = Emu(360000); gap = Emu(200000)
for i, (n, d) in enumerate(viz):
    col = i % 2; row = i // 2
    x = x0 + col * (w + gap)
    y = y0 + row * (h + Emu(100000))
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    round_rect(s, x + Emu(120000), y + Emu(80000),
               Emu(200000), Emu(200000), ACCENT, radius=0.5)
    text(s, x + Emu(400000), y, Emu(1800000), h,
         n, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Emu(2200000), y, w - Emu(2400000), h,
         d, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 30 — Live Demo Walkthrough
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 29, "Part VI · System")
title_block(s, "Live Demo",
            "Six-step walkthrough · ~3 minutes · data seed 2007007")

demo = [
    ("1", "Load synthetic data",
     "Click 'Load Synthetic Data' — deterministic well log loads instantly."),
    ("2", "Preview logs",
     "Data table + 5 log tracks (GR, Rt, RHOB, NPHI, DT) via Recharts."),
    ("3", "Run XGBoost",
     "Zone predictions + productivity chart + feature-importance bars."),
    ("4", "Configure & run GA",
     "Adjust waypoints / generations, then 'Run GA Optimization'."),
    ("5", "3D scene",
     "Rotate, pan, zoom — formation layers and optimized trajectory."),
    ("6", "Survey table & history",
     "Show station table, back-navigate, inspect dashboard run history."),
]
x0 = Emu(1100000); y = Emu(1300000)
w = Emu(10400000); h = Emu(620000); gap = Emu(140000)
for n, name, desc in demo:
    round_rect(s, x0, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, x0, y, Emu(80000), h, ACCENT)
    round_rect(s, x0 + Emu(220000), y + Emu(140000),
               Emu(340000), Emu(340000), ACCENT, radius=0.5)
    text(s, x0 + Emu(220000), y + Emu(140000),
         Emu(340000), Emu(340000), n,
         size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + Emu(700000), y + Emu(100000),
         Emu(2800000), Emu(220000), name,
         size=12, bold=True, color=INK)
    text(s, x0 + Emu(700000), y + Emu(330000),
         w - Emu(900000), Emu(260000), desc,
         size=10, color=MUTED)
    y += h + gap


# =====================================================================
# SLIDE 31 — Key Parameters & Results
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 30, "Part VII · Results")
title_block(s, "Key Parameters & Results",
            "Everything you need to reproduce the demo")

params = [
    ("1000 – 2000 m", "Survey depth range"),
    ("200",           "Data points (5 m sampling)"),
    ("150",           "XGBoost trees (max_depth=5)"),
    ("50",            "GA population size"),
    ("100",           "GA generations (50 – 200)"),
    ("8",             "Waypoints (survey stations)"),
    ("14 genes",      "Chromosome length (2 × 7)"),
    ("0.3",           "DLS weight in fitness"),
    ("5",             "Input log features"),
    ("3",             "Output zone classes"),
    ("2007007",       "RNG seed (Student ID)"),
    ("< 30 s",        "End-to-end runtime on laptop"),
]
cols_n = 4
x0 = Emu(1100000); y0 = Emu(1400000)
w = Emu(2550000); h = Emu(1100000); gx = Emu(120000); gy = Emu(160000)
for i, (val, name) in enumerate(params):
    col = i % cols_n; row = i // cols_n
    x = x0 + col * (w + gx)
    y = y0 + row * (h + gy)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.1)
    rect(s, x, y, w, Emu(70000), ACCENT)
    text(s, x, y + Emu(240000), w, Emu(380000),
         val, size=20, bold=True, color=INK,
         font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, x, y + Emu(680000), w, Emu(280000),
         name, size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 32 — Limitations & Future Work
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 31, "Part VIII · Outlook")
title_block(s, "Limitations & Future Work",
            "Each limitation is a clear extension path toward MSc / PhD")

cols = [("Limitation", 2800000), ("Current state", 4000000),
        ("Extension path", 3600000)]
rows = [
    ("Heuristic labels", "No DST/MDT ground truth",        "Real formation test data"),
    ("Single-well",      "No multi-well correlation",       "Multi-well geostatistics"),
    ("2D formations",    "Horizontal layers, no dip/fault", "Structural geology integration"),
    ("No mechanics",     "No torque / drag / stability",    "Mechanical Earth Model"),
    ("Synthetic data",   "Random formations in demo",       "Real Bangladesh field data"),
]
x0 = Emu(1100000); y = Emu(1300000); row_h = Emu(440000)
cx = x0
for name, w in cols:
    rect(s, cx, y, Emu(w), row_h, ACCENT)
    text(s, cx + Emu(160000), y, Emu(w - 300000), row_h,
         name, size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         anchor=MSO_ANCHOR.MIDDLE)
    cx += Emu(w)
y += row_h
for ri, row in enumerate(rows):
    cx = x0
    bgc = CARD if ri % 2 == 0 else SOFT
    for ci, ((_, w), val) in enumerate(zip(cols, row)):
        rect(s, cx, y, Emu(w), row_h, bgc, line=BORDER, line_w=Emu(3175))
        col_ink = INK if ci == 0 else (INK_2 if ci == 1 else ACCENT_D)
        text(s, cx + Emu(160000), y, Emu(w - 300000), row_h,
             val, size=11, bold=(ci in (0, 2)), color=col_ink,
             anchor=MSO_ANCHOR.MIDDLE)
        cx += Emu(w)
    y += row_h

callout(s, Emu(5900000),
        "Architecture is designed to absorb each upgrade without rewriting the core",
        size=11)


# =====================================================================
# SLIDE 33 — Conclusion
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 32, "Part VIII · Outlook")
title_block(s, "Conclusion",
            "What WellPath.AI demonstrates — the methodological contribution")

round_rect(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(600000),
           ACCENT_L, line=ACCENT_LB, line_w=Emu(9525), radius=0.2)
text(s, Emu(1100000), Emu(1300000), Emu(10400000), Emu(600000),
     "ML + evolutionary computation can automate directional well planning end-to-end",
     size=15, bold=True, color=ACCENT_D,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

points = [
    ("XGBoost",           "Classifies reservoir quality from 5 log curves with interpretable feature importance"),
    ("Genetic Algorithm", "Evolves trajectories that maximize productive contact while respecting DLS"),
    ("Minimum Curvature", "Converts GA angles to 3D coordinates — industry-standard geometry"),
    ("Three.js viewer",   "Interactive 3D output + exportable survey station table for engineers"),
    ("Contribution",      "A complete, open, end-to-end pipeline — methodological, not incremental"),
    ("Self-contained",    "Heuristic labelling requires no external training data; seed 2007007"),
]
x0 = Emu(1100000); y0 = Emu(2100000)
w = Emu(5100000); h = Emu(600000); gx = Emu(200000); gy = Emu(140000)
for i, (n, d) in enumerate(points):
    col = i % 2; row = i // 2
    x = x0 + col * (w + gx)
    y = y0 + row * (h + gy)
    round_rect(s, x, y, w, h, CARD, line=BORDER, line_w=Emu(6350), radius=0.15)
    rect(s, x, y, Emu(80000), h, ACCENT)
    text(s, x + Emu(200000), y + Emu(80000), w - Emu(300000),
         Emu(220000), n, size=12, bold=True, color=INK)
    text(s, x + Emu(200000), y + Emu(300000), w - Emu(300000),
         Emu(280000), d, size=10, color=MUTED)


# =====================================================================
# SLIDE 34 — Thank You
# =====================================================================
s = prs.slides.add_slide(blank)
set_bg(s, BG)
rect(s, 0, 0, Emu(120000), SH, ACCENT)

text(s, Emu(480000), Emu(1600000), Emu(11000000), Emu(1000000),
     "Thank You", size=84, bold=True, color=INK,
     align=PP_ALIGN.CENTER)

text(s, Emu(480000), Emu(2700000), Emu(11000000), Emu(440000),
     "I welcome your questions.", size=20, color=MUTED,
     align=PP_ALIGN.CENTER, italic=True)

rect(s, Emu(5700000), Emu(3400000), Emu(800000), Emu(40000), ACCENT)

text(s, Emu(480000), Emu(3520000), Emu(11000000), Emu(340000),
     "Joseph Ahmed  ·  Student ID 2007007",
     size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, Emu(480000), Emu(3860000), Emu(11000000), Emu(300000),
     "BSc Petroleum & Mining Engineering", size=13, color=INK_2,
     align=PP_ALIGN.CENTER)
text(s, Emu(480000), Emu(4160000), Emu(11000000), Emu(300000),
     "Supervisor: Assistant Professor Aqif Hosain Khan",
     size=13, color=INK_2, align=PP_ALIGN.CENTER)
text(s, Emu(480000), Emu(4460000), Emu(11000000), Emu(300000),
     "Chittagong University of Engineering & Technology  ·  December 2025",
     size=13, color=INK_2, align=PP_ALIGN.CENTER)

text(s, Emu(480000), Emu(5400000), Emu(11000000), Emu(320000),
     "github.com/ahmedjoseph07/WellPath.AI",
     size=13, color=ACCENT_D, font=FONT_MONO, align=PP_ALIGN.CENTER)
text(s, Emu(480000), Emu(5720000), Emu(11000000), Emu(300000),
     "ahmedjoseph11@gmail.com",
     size=12, color=MUTED, align=PP_ALIGN.CENTER)


out = "WellPath_AI_Thesis_Presentation.pptx"
prs.save(out)
print(f"Saved {out}  ·  {len(prs.slides)} slides")
